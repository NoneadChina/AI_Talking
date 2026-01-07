# -*- coding: utf-8 -*-
"""
文件解析模块
支持解析Word、Excel、PowerPoint、PDF、Markdown、HTML等常用文件格式为Markdown文本
"""

import os
import logging
from typing import Optional
from abc import ABC, abstractmethod

logger = logging.getLogger(__name__)


class FileParser(ABC):
    """文件解析器基类"""

    @abstractmethod
    def parse(self, file_path: str) -> str:
        """解析文件并返回Markdown文本"""
        pass

    @abstractmethod
    def supports(self, file_path: str) -> bool:
        """检查是否支持该文件格式"""
        pass


class DocxParser(FileParser):
    """Word文档解析器"""

    def supports(self, file_path: str) -> bool:
        return file_path.lower().endswith(('.docx', '.doc'))

    def parse(self, file_path: str) -> str:
        try:
            from docx import Document
            doc = Document(file_path)
            content = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # 检查段落是否有标题样式
                    if para.style.name.startswith('Heading'):
                        level = int(para.style.name.split(' ')[1]) if len(para.style.name.split(' ')) > 1 else 1
                        content.append('#' * level + ' ' + text)
                    else:
                        content.append(text)
            return '\n\n'.join(content)
        except ImportError:
            logger.warning("python-docx 未安装，无法解析Word文档")
            return f"[无法解析Word文档: {os.path.basename(file_path)} - 缺少python-docx库]"
        except Exception as e:
            logger.error(f"解析Word文档失败: {e}")
            return f"[解析Word文档失败: {os.path.basename(file_path)}]"


class ExcelParser(FileParser):
    """Excel文件解析器"""

    def supports(self, file_path: str) -> bool:
        return file_path.lower().endswith(('.xlsx', '.xls'))

    def parse(self, file_path: str) -> str:
        try:
            import pandas as pd
            content = []
            # 读取所有sheet
            excel_file = pd.ExcelFile(file_path)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                content.append(f"## Sheet: {sheet_name}")
                # 转换为Markdown表格
                content.append(df.to_markdown(index=False))
                content.append("")
            return '\n'.join(content)
        except ImportError:
            logger.warning("pandas 未安装，无法解析Excel文件")
            return f"[无法解析Excel文件: {os.path.basename(file_path)} - 缺少pandas库]"
        except Exception as e:
            logger.error(f"解析Excel文件失败: {e}")
            return f"[解析Excel文件失败: {os.path.basename(file_path)}]"


class PowerPointParser(FileParser):
    """PowerPoint文件解析器"""

    def supports(self, file_path: str) -> bool:
        return file_path.lower().endswith(('.pptx', '.ppt'))

    def parse(self, file_path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content = []
            for i, slide in enumerate(prs.slides):
                slide_content = [f"## Slide {i + 1}"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                if len(slide_content) > 1:
                    content.append('\n'.join(slide_content))
                    content.append("")
            return '\n'.join(content)
        except ImportError:
            logger.warning("python-pptx 未安装，无法解析PowerPoint文件")
            return f"[无法解析PowerPoint文件: {os.path.basename(file_path)} - 缺少python-pptx库]"
        except Exception as e:
            logger.error(f"解析PowerPoint文件失败: {e}")
            return f"[解析PowerPoint文件失败: {os.path.basename(file_path)}]"


class PdfParser(FileParser):
    """PDF文件解析器"""

    def supports(self, file_path: str) -> bool:
        return file_path.lower().endswith('.pdf')

    def parse(self, file_path: str) -> str:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            content = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text and text.strip():
                    content.append(f"## Page {i + 1}")
                    content.append(text.strip())
                    content.append("")
            return '\n'.join(content)
        except ImportError:
            logger.warning("PyPDF2 未安装，无法解析PDF文件")
            return f"[无法解析PDF文件: {os.path.basename(file_path)} - 缺少PyPDF2库]"
        except Exception as e:
            logger.error(f"解析PDF文件失败: {e}")
            return f"[解析PDF文件失败: {os.path.basename(file_path)}]"


class MarkdownParser(FileParser):
    """Markdown文件解析器"""

    def supports(self, file_path: str) -> bool:
        return file_path.lower().endswith('.md')

    def parse(self, file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"解析Markdown文件失败: {e}")
            return f"[解析Markdown文件失败: {os.path.basename(file_path)}]"


class HtmlParser(FileParser):
    """HTML文件解析器"""

    def supports(self, file_path: str) -> bool:
        return file_path.lower().endswith('.html')

    def parse(self, file_path: str) -> str:
        try:
            from html import unescape
            from bs4 import BeautifulSoup

            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, 'html.parser')
            content = []

            # 提取标题
            title = soup.find('title')
            if title:
                content.append(f"# {title.get_text().strip()}")

            # 提取主要文本内容
            for tag in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li', 'pre', 'code']):
                text = tag.get_text().strip()
                if text and len(text) > 0:
                    if tag.name.startswith('h'):
                        level = int(tag.name[1])
                        content.append('#' * level + ' ' + text)
                    elif tag.name == 'p':
                        content.append(text)
                    elif tag.name == 'li':
                        content.append('- ' + text)
                    elif tag.name in ['pre', 'code']:
                        content.append(f'`{text}`')

            return '\n\n'.join(content)
        except ImportError:
            logger.warning("beautifulsoup4 未安装，无法解析HTML文件")
            return f"[无法解析HTML文件: {os.path.basename(file_path)} - 缺少beautifulsoup4库]"
        except Exception as e:
            logger.error(f"解析HTML文件失败: {e}")
            return f"[解析HTML文件失败: {os.path.basename(file_path)}]"


class TextParser(FileParser):
    """纯文本文件解析器"""

    def supports(self, file_path: str) -> bool:
        return file_path.lower().endswith('.txt')

    def parse(self, file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"解析文本文件失败: {e}")
            return f"[解析文本文件失败: {os.path.basename(file_path)}]"


class FileParserManager:
    """文件解析管理器"""

    def __init__(self):
        self.parsers = [
            DocxParser(),
            ExcelParser(),
            PowerPointParser(),
            PdfParser(),
            MarkdownParser(),
            HtmlParser(),
            TextParser(),
        ]

    def parse_file(self, file_path: str) -> str:
        """解析文件并返回Markdown文本"""
        if not os.path.exists(file_path):
            logger.error(f"文件不存在: {file_path}")
            return f"[文件不存在: {os.path.basename(file_path)}]"

        for parser in self.parsers:
            if parser.supports(file_path):
                logger.info(f"使用 {parser.__class__.__name__} 解析文件: {file_path}")
                return parser.parse(file_path)

        # 不支持的文件格式
        logger.warning(f"不支持的文件格式: {file_path}")
        return f"[不支持的文件格式: {os.path.basename(file_path)}]"

    def get_supported_extensions(self) -> list:
        """获取支持的文件扩展名"""
        extensions = []
        for parser in self.parsers:
            if hasattr(parser, 'supports'):
                # 从supports方法中提取扩展名
                pass
        return extensions


# 全局文件解析管理器实例
file_parser_manager = FileParserManager()


def parse_file(file_path: str) -> str:
    """解析文件并返回Markdown文本（便捷函数）"""
    return file_parser_manager.parse_file(file_path)
